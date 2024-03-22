import fitz  # PyMuPDF
import os
from ftfy import fix_text
import csv
import docx
from docx.shared import Pt
from pptx import Presentation


from PIL import Image
from pdf2image import convert_from_path
import base64
import pandas as pd
import tabula



# TODO: Add support for multiple formats, e.g. docx, pptx,txt etc. --> done
# TODO: Add support images and tables extraction
# for pdf - done 
# for docx - not working 
# images inside metadata

# TODO: Add support for multiple languages, and mathamatical equations
# TODO: page chunking -> further chunks block wise (in memory vector store using faiss) 

class PDFToCSVConverter:
    def __init__(self, pdf_path, csv_file_name):
        self.pdf_path = pdf_path
        self.csv_file_name = csv_file_name
        self.pdf_document = None

    def open_pdf(self):
        self.pdf_document = fitz.open(self.pdf_path)

    def close_pdf(self):
        if self.pdf_document:
            self.pdf_document.close()

    def extract_text_from_page(self, page_number):
        page = self.pdf_document.load_page(page_number)
        text = page.get_text()
        return text

    def extract_text_from_pages_by_blocks(self, page_number):
        page = self.pdf_document.load_page(page_number)
        text = ""
        for block in page.get_text("blocks"):
            text += block[4] + "\n"
            # also extract images 

        # fix text
        text = fix_text(text)
        return text

    def extract_text_from_all_pages(self, mode="page"):
        extracted_data = {}
        for page_number in range(self.pdf_document.page_count):
            if mode == "block":
                text = self.extract_text_from_pages_by_blocks(page_number)
            else:
                text = self.extract_text_from_page(page_number)
            extracted_data[page_number + 1] = text
        return extracted_data

    def extract_images_and_update_csv(self):
        pdf_document = fitz.open(self.pdf_path)
        output_csv = "output.csv"  # Update this with your original output CSV file

        rows = []
        df = pd.read_csv(output_csv)
        fieldnames = list(df.columns) + ['Images (Base64)']

        for idx, row in df.iterrows():
            page_number = int(row['Page Number']) - 1
            page = pdf_document.load_page(page_number)
            image_list = page.get_images(full=True)

            images_base64 = []
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                images_base64.append(image_base64)

            row['Images (Base64)'] = images_base64
            rows.append(row.to_dict())

        df = pd.DataFrame(rows, columns=fieldnames)
        df.to_csv(output_csv, index=False)

        pdf_document.close()

    def extract_tables_from_page(self, page_number):
        tables = tabula.read_pdf(self.pdf_path, pages=page_number+1, multiple_tables=True)
        formatted_tables = []

        for table in tables:
            table_data = table.values.tolist()
            formatted_table = '\n'.join(['\t'.join(map(str, row)) for row in table_data])
            formatted_table = '\n'.join(filter(None, map(str.strip, formatted_table.splitlines())))
            formatted_tables.append(formatted_table)

        return formatted_tables

    def save_to_csv(self, data):
        pdf_name = os.path.basename(self.pdf_path)
        rows = []
        for page_number, content in data.items():
            rows.append([pdf_name, page_number, content])

        df = pd.DataFrame(rows, columns=["file_name", "Page Number", "Content"])
        df.to_csv(self.csv_file_name, index=False)

    def convert(self, mode="page"):
        self.open_pdf()
        data = self.extract_text_from_all_pages(mode)
        self.save_to_csv(data)
        self.close_pdf()



class WordToCSVConverter:
    def __init__(self, docx_path, csv_file_name="output.csv", paragraphs_per_page=10):
        self.docx_path = docx_path
        self.csv_file_name = csv_file_name
        self.paragraphs_per_page = paragraphs_per_page
        self.docx_document = None

    def open_docx(self):
        self.docx_document = docx.Document(self.docx_path)

    def close_docx(self):
        self.docx_document = None

    def extract_text_page_wise(self):
        text = ""
        current_page = 1
        paragraph_count = 0

        for paragraph in self.docx_document.paragraphs:
            text += paragraph.text + "\n"
            paragraph_count += 1

            if paragraph_count >= self.paragraphs_per_page:
                text += f"\n[Page {current_page + 1}]\n"
                current_page += 1
                paragraph_count = 0

        text = self.fix_text(text)
        return text

    def extract_images_page_wise(self):
        images = []
        current_page = 1

        for rel in self.docx_document.part.rels:
            if "image" in str(self.docx_document.part.rels[rel].target_ref):
                image_blob = self.docx_document.part.rels[rel].target_part.blob
                image_base64 = base64.b64encode(image_blob).decode('utf-8')
                images.append(image_base64)

        return images

    def extract_from_folder(self, folder_path):
        for file in os.listdir(folder_path):
            if file.endswith(".docx"):
                docx_path = os.path.join(folder_path, file)
                self.docx_path = docx_path
                self.open_docx()
                text_data = self.extract_text_page_wise()
                image_data = self.extract_images_page_wise()
                self.save_to_csv(text_data, image_data)
                self.close_docx()

    def save_to_csv(self, text_data, image_data):
        docx_name = os.path.basename(self.docx_path)
        rows = []
        current_page = 1

        for content in text_data.split('\n'):
            if content.startswith("[Page"):
                current_page = int(content.split()[1][:-1])
            else:
                images_on_page = []
                if current_page <= len(image_data):
                    images_on_page = image_data[current_page - 1]

                rows.append([docx_name, current_page, content, images_on_page])

        df = pd.DataFrame(rows, columns=["Doc Name", "Page Number", "Content", "Image Data"])
        df.to_csv(self.csv_file_name, index=False)

    @staticmethod
    def fix_text(text):
        # text fixing logic here
        return text


class PPTXToCSVConverter:
    def __init__(self, pptx_path, csv_file_name):
        self.pptx_path = pptx_path
        self.csv_file_name = csv_file_name
        self.pptx_presentation = None

    def open_pptx(self):
        self.pptx_presentation = Presentation(self.pptx_path)

    def close_pptx(self):
        self.pptx_presentation = None  # PowerPoint presentations don't need closing

    def extract_text_from_slides(self):
        text = ""
        for slide in self.pptx_presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                    # Add image extraction logic if needed

        text = self.fix_text(text)
        return text

    def extract_text_from_all_slides(self):
        extracted_data = {}
        for idx, slide in enumerate(self.pptx_presentation.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                    # Add image extraction logic if needed

            text = self.fix_text(text)
            extracted_data[idx + 1] = text
        return extracted_data
 
    def extract_from_folder(self, folder_path):
        for file in os.listdir(folder_path):
            if file.endswith(".pptx"):
                pptx_path = os.path.join(folder_path, file)
                self.pptx_path = pptx_path
                self.open_pptx()
                data = self.extract_text_from_all_slides()
                self.save_to_csv(data)
                self.close_pptx()

    def convert(self):
        self.open_pptx()
        data = self.extract_text_from_all_slides()
        self.save_to_csv(data)
        self.close_pptx()

    def save_to_csv(self, data):
        pptx_name = os.path.basename(self.pptx_path)
        rows = []
        for slide_number, content in data.items():
            rows.append([pptx_name, slide_number, content])

        df = pd.DataFrame(rows, columns=["file_name", "Slide Number", "Content"])
        df.to_csv(self.csv_file_name, index=False)

    @staticmethod
    def fix_text(text):
        # text fixing logic here
        return text

class TextToCSVConverter:
    def __init__(self, txt_path, csv_file_name):
        self.txt_path = txt_path
        self.csv_file_name = csv_file_name

    def read_text(self):
        with open(self.txt_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return text

    def save_to_csv(self, text):
        txt_name = os.path.basename(self.txt_path)
        df = pd.DataFrame({"file_name": [txt_name], "Content": [text]})
        df.to_csv(self.csv_file_name, index=False)

    def save_transcript_to_csv(self, transcript):
        txt_name = os.path.basename(self.txt_path)
        df = pd.DataFrame({"file_name": [txt_name], "Transcript": [transcript]})
        df.to_csv(self.csv_file_name, index=False, mode='a', header=False)

    def convert(self):
        text_content = self.read_text()
        transcript = get_video_transcript(video_id)  # Replace 'video_id' with the actual YouTube video ID
        self.save_to_csv(text_content)
        self.save_transcript_to_csv(transcript)




# logic to check file type and call class accordingly
def convert_files_in_folder(file_path):

    if file_path.lower().endswith(".pdf"):
        pdf_converter = PDFToCSVConverter(file_path, 'output.csv')
        print("called")
        # pdf_converter.extract_from_folder(folder_path)
        pdf_converter.convert()
        pdf_converter.extract_images_and_update_csv()
    
    elif file_path.lower().endswith(".docx"):
        docx_converter = WordToCSVConverter(file_path, 'output.csv')
        docx_converter.extract_from_folder(folder_path)
        # docx_converter.convert()
    
    elif file_path.lower().endswith(".pptx"):
        pptx_converter = PPTXToCSVConverter(file_path, 'output.csv')
        pptx_converter.extract_from_folder(folder_path)

    elif file_path.lower().endswith(".txt"):
        txt_converter = TextToCSVConverter(file_path, 'output.csv')
        txt_converter.convert()
    
    else:
        print(f"Unsupported file format: {file}")


# Usage
# folder_path = r"C:\Users\abhin\OneDrive\Desktop\hp glc code\glc-hp-hackathon\chromadb\files_to_extract"  # Replace this with your folder path
# convert_files_in_folder(folder_path)
